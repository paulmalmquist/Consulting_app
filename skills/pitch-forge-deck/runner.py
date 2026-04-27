#!/usr/bin/env python3
"""pitch-forge-deck/runner.py — Standalone runner for the Pitch Forge deck skill.

Three subcommands:
  generate   — call OpenAI to produce structured slide content JSON + markdown outline
  build      — convert slides.json into a .pptx using python-pptx
  critique   — run Sarat Mode critique on slides.json, write critique.json + summary

Usage (from backend/ directory):
  python ../skills/pitch-forge-deck/runner.py generate \
    --topic "AP Automation for Hall Boys" \
    --client "Hall Boys Holdings" \
    --iteration 1 \
    --constraints "No new headcount, Acumatica is system of record" \
    --output-dir "../demo_docs/pitch_forge_deck/iteration_1"

  python ../skills/pitch-forge-deck/runner.py build \
    --slides-file "../demo_docs/pitch_forge_deck/iteration_1/slides.json" \
    --output "../demo_docs/pitch_forge_deck/iteration_1/presentation.pptx" \
    --client "Hall Boys Holdings" \
    --iteration 1

  python ../skills/pitch-forge-deck/runner.py critique \
    --slides-file "../demo_docs/pitch_forge_deck/iteration_1/slides.json" \
    --client "Hall Boys Holdings" \
    --iteration 1 --max-iterations 3 \
    --output-dir "../demo_docs/pitch_forge_deck/iteration_1"
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Bootstrap: backend/ must be on sys.path for app.* imports
# ---------------------------------------------------------------------------

SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parent.parent
BACKEND_DIR = REPO_ROOT / "backend"
if str(BACKEND_DIR) not in sys.path:
    sys.path.insert(0, str(BACKEND_DIR))


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _write_json(path: Path, data) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, indent=2, default=str), encoding="utf-8")
    print(f"  wrote {path}")


def _write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    print(f"  wrote {path}")


def _chat(prompt: str, temperature: float = 0.15) -> str:
    from app.services.ai_client import get_instrumented_sync_client
    client = get_instrumented_sync_client()
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature,
        max_tokens=4096,
    )
    return response.choices[0].message.content.strip()


def _parse_json(raw: str, context: str):
    cleaned = re.sub(r"```(?:json)?\s*", "", raw).replace("```", "").strip()
    for start_char, end_char in [('{', '}'), ('[', ']')]:
        idx = cleaned.find(start_char)
        if idx != -1:
            depth = 0
            for i, ch in enumerate(cleaned[idx:], start=idx):
                if ch == start_char:
                    depth += 1
                elif ch == end_char:
                    depth -= 1
                    if depth == 0:
                        try:
                            return json.loads(cleaned[idx:i+1])
                        except json.JSONDecodeError as exc:
                            raise ValueError(f"JSON parse error in {context}: {exc}") from exc
    raise ValueError(f"No JSON found in {context} response")


def _check_banned(text: str) -> list[str]:
    from app.services.pitch_forge_prompts import BANNED_PHRASES
    lower = text.lower()
    return [p for p in BANNED_PHRASES if p.lower() in lower]


# ---------------------------------------------------------------------------
# Subcommand: generate
# ---------------------------------------------------------------------------

SLIDE_GENERATION_PROMPT = """You are building a 12-slide executive presentation for {client_name}.

Topic: {topic}
Iteration: {iteration_num} of {max_iterations}
Known constraints: {constraints}

{previous_critique_block}

Generate exactly 12 slides. Return ONLY a JSON object in this exact format:

{{
  "title": "Presentation title",
  "subtitle": "One-line subtitle",
  "slides": [
    {{
      "slide_num": 1,
      "layout": "title",
      "title": "...",
      "subtitle": "...",
      "notes": "Speaker note here"
    }},
    {{
      "slide_num": 2,
      "layout": "section_header",
      "title": "Section title",
      "body": "One sentence framing",
      "notes": "..."
    }},
    {{
      "slide_num": 3,
      "layout": "two_column",
      "title": "Slide title",
      "left_header": "Left column header",
      "left_points": ["Point 1", "Point 2", "Point 3"],
      "right_header": "Right column header",
      "right_points": ["Point 1", "Point 2", "Point 3"],
      "notes": "..."
    }},
    ...
  ]
}}

Layout options: title | section_header | bullets | two_column | stat_callout | quote | table | closing

Rules:
- Slide 1: always "title" layout
- Slide 12: always "closing" layout
- Every claim must tie to {client_name} specifically — no generic industry statements
- Impact claims must include a number or timeframe (e.g. "3 hours saved per invoice cycle" not "significant time savings")
- No filler phrases: "streamline", "leverage", "unlock", "synergy", "transform", "revolutionize", "game-changing", "cutting-edge", "world-class", "best-in-class", "seamlessly", "holistic", "end-to-end", "robust", "scalable solution", "empower", "optimize", "innovative"
- stat_callout layout has: title, stat (large number), stat_label, body
- quote layout has: title, quote, attribution
- table layout has: title, headers (list), rows (list of lists)
- bullets layout has: title, points (list of strings), optional subtitle
"""


def cmd_generate(args: argparse.Namespace) -> int:
    from app.services.pitch_forge_prompts import BANNED_PHRASES

    out_dir = Path(args.output_dir)

    # Load previous critique if provided
    previous_critique_block = ""
    if args.critique_file:
        cp = Path(args.critique_file)
        if cp.exists():
            critique = json.loads(cp.read_text())
            overall = critique.get("overall_verdict", "REJECT")
            summary = critique.get("sarat_voice_summary", "")
            kills = critique.get("kill_list", [])
            fixable = critique.get("fixable_issues", [])
            kill_text = "\n".join(f"  - [{k.get('lens','')}] {k.get('section','')}: {k.get('reason','')}" for k in kills)
            fix_text = "\n".join(f"  - {f}" for f in fixable)
            previous_critique_block = f"""Previous iteration verdict: {overall}
Sarat said: "{summary}"
Killed sections:
{kill_text or '  (none)'}
Address before rebuilding:
{fix_text or '  (none)'}

In this iteration, directly fix the items Sarat killed or flagged. Do not repeat the same mistakes."""

    constraints_text = args.constraints or "None specified"

    prompt = SLIDE_GENERATION_PROMPT.format(
        client_name=args.client,
        topic=args.topic,
        iteration_num=args.iteration,
        max_iterations=args.max_iterations,
        constraints=constraints_text,
        previous_critique_block=previous_critique_block,
    )

    print(f"  Generating slide content (iteration {args.iteration})...")
    raw = _chat(prompt, temperature=0.15)
    data = _parse_json(raw, "slide_generation")

    # Check banned phrases across all text
    all_text = json.dumps(data)
    found = _check_banned(all_text)
    if found:
        print(f"  WARNING: banned phrases found: {found} — retrying once")
        raw = _chat(prompt + f"\n\nDo NOT use these phrases: {found}", temperature=0.1)
        data = _parse_json(raw, "slide_generation_retry")
        found = _check_banned(json.dumps(data))
        if found:
            print(f"  ERROR: banned phrases persist after retry: {found}")
            return 1

    _write_json(out_dir / "slides.json", data)

    # Write markdown outline
    lines = [f"# {data.get('title', 'Presentation')}", f"*{data.get('subtitle', '')}*", ""]
    for slide in data.get("slides", []):
        lines.append(f"## Slide {slide['slide_num']}: {slide.get('title', '')}")
        layout = slide.get("layout", "")
        if layout == "bullets":
            for pt in slide.get("points", []):
                lines.append(f"- {pt}")
        elif layout == "two_column":
            lines.append(f"**{slide.get('left_header','')}**")
            for pt in slide.get("left_points", []):
                lines.append(f"- {pt}")
            lines.append(f"**{slide.get('right_header','')}**")
            for pt in slide.get("right_points", []):
                lines.append(f"- {pt}")
        elif layout == "stat_callout":
            lines.append(f"**{slide.get('stat','')}** {slide.get('stat_label','')}")
            lines.append(slide.get("body", ""))
        elif layout == "quote":
            lines.append(f"> {slide.get('quote','')}")
            lines.append(f"— {slide.get('attribution','')}")
        else:
            if slide.get("body"):
                lines.append(slide["body"])
        lines.append("")

    _write_text(out_dir / "presentation.md", "\n".join(lines))
    print(f"  Generated {len(data.get('slides', []))} slides.")
    return 0


# ---------------------------------------------------------------------------
# Subcommand: build
# ---------------------------------------------------------------------------

def cmd_build(args: argparse.Namespace) -> int:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx --break-system-packages")
        return 1

    slides_path = Path(args.slides_file)
    if not slides_path.exists():
        print(f"ERROR: slides file not found: {slides_path}")
        return 1

    data = json.loads(slides_path.read_text())
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Design tokens — Midnight Executive palette
    DARK_BG  = RGBColor(0x1E, 0x27, 0x61)   # navy
    MID_BG   = RGBColor(0x16, 0x1B, 0x4A)   # deeper navy for section headers
    LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)    # white for content slides
    ACCENT   = RGBColor(0xCA, 0xDC, 0xFC)   # ice blue
    TEXT_DARK = RGBColor(0xFF, 0xFF, 0xFF)   # white on dark
    TEXT_LIGHT = RGBColor(0x1E, 0x27, 0x61)  # navy on white
    TEXT_MUTED = RGBColor(0x55, 0x65, 0xA0)  # muted blue-grey
    HIGHLIGHT = RGBColor(0x4A, 0x9E, 0xFF)   # bright blue accent

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_rect(slide, x, y, w, h, color):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, font_size, color, bold=False,
                 italic=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
        return txBox

    def add_bullet_box(slide, points, x, y, w, h, font_size, color, header=None):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        if header:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = header
            run.font.size = Pt(font_size + 2)
            run.font.bold = True
            run.font.color.rgb = ACCENT
            run.font.name = "Calibri"
            first = False
        for pt in points:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            run = p.add_run()
            run.text = f"• {pt}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
            first = False

    print(f"  Building PPTX ({len(data.get('slides', []))} slides)...")

    for slide_data in data.get("slides", []):
        layout = slide_data.get("layout", "bullets")
        slide = prs.slides.add_slide(blank_layout)
        title = slide_data.get("title", "")
        slide_num = slide_data.get("slide_num", 0)

        if layout == "title":
            # Full dark background
            add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
            # Accent bar at bottom third
            add_rect(slide, 0, 5.0, 13.33, 0.08, ACCENT)
            # Title
            add_text(slide, data.get("title", title), 1.0, 1.8, 11.33, 1.8,
                     44, TEXT_DARK, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
            # Subtitle
            sub = slide_data.get("subtitle", data.get("subtitle", ""))
            if sub:
                add_text(slide, sub, 1.5, 3.8, 10.33, 0.8,
                         20, ACCENT, italic=True, align=PP_ALIGN.CENTER)
            # Client / iteration footnote
            add_text(slide, f"{args.client}  •  Iteration {args.iteration}",
                     1.5, 6.6, 10.33, 0.5, 12, TEXT_MUTED, align=PP_ALIGN.CENTER)

        elif layout == "section_header":
            add_rect(slide, 0, 0, 13.33, 7.5, MID_BG)
            add_rect(slide, 0.6, 3.4, 0.08, 0.7, ACCENT)
            add_text(slide, title, 1.0, 3.0, 11.33, 1.2,
                     36, TEXT_DARK, bold=True, font_name="Georgia")
            body = slide_data.get("body", "")
            if body:
                add_text(slide, body, 1.0, 4.3, 10.33, 0.8, 16, ACCENT, italic=True)

        elif layout == "closing":
            add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
            add_rect(slide, 0, 3.5, 13.33, 0.08, ACCENT)
            add_text(slide, title, 1.0, 1.5, 11.33, 1.4,
                     40, TEXT_DARK, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
            body = slide_data.get("body", slide_data.get("subtitle", ""))
            if body:
                add_text(slide, body, 1.5, 3.2, 10.33, 1.2, 18, ACCENT,
                         italic=True, align=PP_ALIGN.CENTER)
            add_text(slide, args.client, 1.5, 6.6, 10.33, 0.5,
                     12, TEXT_MUTED, align=PP_ALIGN.CENTER)

        elif layout == "stat_callout":
            add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG)
            add_rect(slide, 0, 0, 13.33, 1.1, DARK_BG)
            add_text(slide, title, 0.5, 0.15, 12.33, 0.8, 24, TEXT_DARK, bold=True)
            # Big stat
            stat = slide_data.get("stat", "")
            stat_label = slide_data.get("stat_label", "")
            add_text(slide, stat, 1.5, 1.5, 10.33, 2.2,
                     72, DARK_BG, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
            if stat_label:
                add_text(slide, stat_label, 1.5, 3.7, 10.33, 0.6,
                         18, TEXT_MUTED, align=PP_ALIGN.CENTER)
            body = slide_data.get("body", "")
            if body:
                add_text(slide, body, 1.0, 4.6, 11.33, 1.5, 14, TEXT_LIGHT)

        elif layout == "two_column":
            add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG)
            add_rect(slide, 0, 0, 13.33, 1.1, DARK_BG)
            add_text(slide, title, 0.5, 0.15, 12.33, 0.8, 24, TEXT_DARK, bold=True)
            # Divider
            add_rect(slide, 6.565, 1.3, 0.05, 5.8, RGBColor(0xCC, 0xD5, 0xE8))
            add_bullet_box(slide,
                           slide_data.get("left_points", []),
                           0.5, 1.3, 5.8, 5.8, 14, TEXT_LIGHT,
                           header=slide_data.get("left_header"))
            add_bullet_box(slide,
                           slide_data.get("right_points", []),
                           6.83, 1.3, 5.8, 5.8, 14, TEXT_LIGHT,
                           header=slide_data.get("right_header"))

        elif layout == "quote":
            add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
            add_text(slide, "“", 0.8, 1.0, 1.5, 2.0, 120, ACCENT, bold=True)
            add_text(slide, slide_data.get("quote", ""), 1.5, 1.5, 10.33, 3.0,
                     22, TEXT_DARK, italic=True, align=PP_ALIGN.CENTER)
            add_rect(slide, 4.5, 4.9, 4.33, 0.05, ACCENT)
            add_text(slide, slide_data.get("attribution", ""), 1.5, 5.1, 10.33, 0.6,
                     14, ACCENT, align=PP_ALIGN.CENTER)

        elif layout == "table":
            add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG)
            add_rect(slide, 0, 0, 13.33, 1.1, DARK_BG)
            add_text(slide, title, 0.5, 0.15, 12.33, 0.8, 24, TEXT_DARK, bold=True)
            headers = slide_data.get("headers", [])
            rows = slide_data.get("rows", [])
            if headers and rows:
                cols = len(headers)
                col_w = 12.0 / cols
                # Header row
                for ci, h in enumerate(headers):
                    add_rect(slide, 0.5 + ci * col_w, 1.2, col_w - 0.05, 0.45, DARK_BG)
                    add_text(slide, h, 0.55 + ci * col_w, 1.25, col_w - 0.1, 0.38,
                             12, TEXT_DARK, bold=True)
                for ri, row in enumerate(rows[:6]):
                    bg = RGBColor(0xF0, 0xF4, 0xFB) if ri % 2 == 0 else LIGHT_BG
                    for ci, cell in enumerate(row[:cols]):
                        add_rect(slide, 0.5 + ci * col_w, 1.7 + ri * 0.55,
                                 col_w - 0.05, 0.5, bg)
                        add_text(slide, str(cell),
                                 0.55 + ci * col_w, 1.72 + ri * 0.55,
                                 col_w - 0.1, 0.44, 11, TEXT_LIGHT)

        else:  # bullets (default)
            add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG)
            add_rect(slide, 0, 0, 13.33, 1.1, DARK_BG)
            add_text(slide, title, 0.5, 0.15, 12.33, 0.8, 24, TEXT_DARK, bold=True)
            subtitle = slide_data.get("subtitle", "")
            if subtitle:
                add_text(slide, subtitle, 0.6, 1.2, 12.13, 0.55, 14, TEXT_MUTED, italic=True)
            points = slide_data.get("points", [])
            if points:
                y_start = 1.85 if subtitle else 1.3
                add_bullet_box(slide, points, 0.7, y_start, 12.0, 5.5, 15, TEXT_LIGHT)

        # Slide number (bottom right, all slides except title/closing)
        if layout not in ("title", "closing"):
            add_text(slide, str(slide_num), 12.5, 7.1, 0.6, 0.3,
                     10, TEXT_MUTED, align=PP_ALIGN.RIGHT)

    prs.save(str(output_path))
    print(f"  Saved: {output_path}")
    return 0


# ---------------------------------------------------------------------------
# Subcommand: critique
# ---------------------------------------------------------------------------

def cmd_critique(args: argparse.Namespace) -> int:
    from app.services.pitch_forge_prompts import (
        sarat_mode_critique_prompt,
        SARAT_OBJECTION_STACK,
    )

    slides_path = Path(args.slides_file)
    if not slides_path.exists():
        print(f"ERROR: slides file not found: {slides_path}")
        return 1

    out_dir = Path(args.output_dir)
    slides_data = json.loads(slides_path.read_text())
    constraints_list = [{"text": c.strip()} for c in (args.constraints or "").split(",") if c.strip()]

    # Convert slides JSON to a text summary for the critique prompt
    lines = [f"Presentation: {slides_data.get('title', 'Untitled')}", ""]
    for s in slides_data.get("slides", []):
        lines.append(f"Slide {s['slide_num']} [{s.get('layout','')}]: {s.get('title','')}")
        for field in ("subtitle", "body", "stat", "stat_label", "quote"):
            if s.get(field):
                lines.append(f"  {field}: {s[field]}")
        for field in ("points", "left_points", "right_points"):
            for pt in s.get(field, []):
                lines.append(f"  • {pt}")
        if s.get("rows"):
            for row in s["rows"]:
                lines.append(f"  row: {row}")
    presentation_text = "\n".join(lines)

    print(f"  Running Sarat Mode critique (iteration {args.iteration})...")
    prompt = sarat_mode_critique_prompt(
        client_name=args.client,
        presentation_content=presentation_text,
        iteration_num=args.iteration,
        max_iterations=args.max_iterations,
        known_constraints=constraints_list,
    )

    raw = _chat(prompt, temperature=0.2)

    try:
        review = _parse_json(raw, "sarat_critique")
    except ValueError as exc:
        print(f"  ERROR: could not parse Sarat critique: {exc}")
        return 1

    # Validate
    overall = review.get("overall_verdict")
    if overall not in ("PASS", "WEAK", "REJECT"):
        print(f"  WARNING: unexpected overall_verdict {overall!r} — defaulting to WEAK")
        review["overall_verdict"] = "WEAK"

    score = review.get("score", 0)
    review.setdefault("section_verdicts", {})
    review.setdefault("objection_scores", {})
    review.setdefault("sarat_voice_summary", "No summary provided.")
    review.setdefault("kill_list", [])
    review.setdefault("fatal_issues", [])
    review.setdefault("fixable_issues", [])

    _write_json(out_dir / "critique.json", review)

    # Human-readable summary
    sv = review["section_verdicts"]
    obj = sorted(review["objection_scores"].items(), key=lambda x: -x[1])
    emoji = {"PASS": "✅", "WEAK": "⚠️", "REJECT": "❌"}
    lines = [
        f"# Sarat Critique — Iteration {args.iteration}",
        f"**Overall: {emoji.get(overall,'')} {overall} — Score {score}/100**",
        "",
        f"> {review['sarat_voice_summary']}",
        "",
        "## Section Verdicts",
    ]
    for section, verdict in sv.items():
        lines.append(f"- {emoji.get(verdict,'')} **{verdict}** — {section}")

    if obj:
        lines += ["", "## Objection Scores"]
        for criterion, s in obj:
            bar = "█" * s + "░" * (10 - s)
            lines.append(f"- {criterion}: {bar} {s}/10")

    if review["kill_list"]:
        lines += ["", "## Killed"]
        for k in review["kill_list"]:
            lines.append(f"- **{k.get('section','')}** [{k.get('lens','')}]: {k.get('reason','')}")

    if review["fixable_issues"]:
        lines += ["", "## Fix Before Next Iteration"]
        for f_item in review["fixable_issues"]:
            lines.append(f"- {f_item}")

    _write_text(out_dir / "critique_summary.md", "\n".join(lines))
    print(f"  Verdict: {emoji.get(overall,'')} {overall} — Score {score}/100")
    print(f"  Sarat: {review['sarat_voice_summary'][:120]}...")
    return 0


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(description="Pitch Forge deck skill runner")
    sub = parser.add_subparsers(dest="command", required=True)

    # generate
    g = sub.add_parser("generate", help="Generate slide content JSON")
    g.add_argument("--topic", required=True)
    g.add_argument("--client", default="the client")
    g.add_argument("--iteration", type=int, default=1)
    g.add_argument("--max-iterations", type=int, default=3)
    g.add_argument("--constraints", default="")
    g.add_argument("--critique-file", default=None,
                   help="Path to previous iteration's critique.json (skip on iteration 1)")
    g.add_argument("--output-dir", required=True)

    # build
    b = sub.add_parser("build", help="Build .pptx from slides.json")
    b.add_argument("--slides-file", required=True)
    b.add_argument("--output", required=True)
    b.add_argument("--client", default="the client")
    b.add_argument("--iteration", type=int, default=1)

    # critique
    c = sub.add_parser("critique", help="Run Sarat Mode critique")
    c.add_argument("--slides-file", required=True)
    c.add_argument("--client", default="the client")
    c.add_argument("--iteration", type=int, default=1)
    c.add_argument("--max-iterations", type=int, default=3)
    c.add_argument("--constraints", default="")
    c.add_argument("--output-dir", required=True)

    args = parser.parse_args()

    if args.command == "generate":
        return cmd_generate(args)
    elif args.command == "build":
        return cmd_build(args)
    elif args.command == "critique":
        return cmd_critique(args)
    else:
        parser.print_help()
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
